# -*- coding: utf-8 -*-
"""
Gera o PowerPoint do tutorial teórico do projeto MinhaPrimeiraApi.
Requer: pip install python-pptx
Uso: python gerar_apresentacao.py
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Instale o pacote: pip install python-pptx")
    raise

OUTPUT_FILE = Path(__file__).resolve().parent / "MinhaPrimeiraApi-Tutorial-Aula.pptx"


def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets, notes=""):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[i] if i < len(body.paragraphs) else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "MinhaPrimeiraApi — Tutorial Teórico",
        "Ferramentas: Docker, Entity Framework Core, MySQL\nDesenvolvimento Web com C#"
    )

    add_content_slide(
        prs,
        "Objetivos da aula",
        [
            "Conhecer as ferramentas utilizadas no projeto: Docker, MySQL e Entity Framework Core.",
            "Entender o papel de cada uma na arquitetura da API.",
            "Saber executar a aplicação passo a passo (com e sem Docker).",
            "Ter uma visão geral dos endpoints e da estrutura do projeto."
        ],
        notes="Apresente os objetivos antes de entrar nas ferramentas."
    )

    add_content_slide(
        prs,
        "Visão geral do projeto",
        [
            "API REST em ASP.NET Core (.NET 9) para fins didáticos.",
            "Recursos: Produtos e Categorias (CRUD completo).",
            "Banco de dados MySQL para persistência.",
            "Containerização com Docker: MySQL + API em um único comando.",
            "Entity Framework Core como ORM para acesso ao banco."
        ],
        notes="Projeto usado em sala para ilustrar API + banco + Docker."
    )

    add_content_slide(
        prs,
        "Ferramenta 1: Docker",
        [
            "Docker: plataforma de virtualização em nível de sistema operacional (containers).",
            "Container: ambiente isolado e leve que empacota aplicação + dependências.",
            "Vantagens: mesmo ambiente em dev e produção; fácil reprodução; um comando para subir tudo.",
            "No projeto: uma imagem para o MySQL e uma para a API; o Docker Compose orquestra os dois."
        ],
        notes="Diferencie container de VM: mais leve, compartilha o kernel."
    )

    add_content_slide(
        prs,
        "Conceitos Docker (resumo)",
        [
            "Imagem: modelo read-only (ex.: mysql:8.0, nossa API buildada pelo Dockerfile).",
            "Container: instância em execução de uma imagem.",
            "Dockerfile: receita para construir a imagem da nossa API (build + run em estágios).",
            "Docker Compose: arquivo YAML que define e sobe vários serviços (mysql + api) e a rede entre eles."
        ],
        notes="Mostre o docker-compose.yml do projeto na IDE."
    )

    add_content_slide(
        prs,
        "Ferramenta 2: MySQL",
        [
            "MySQL: banco de dados relacional (SGBD) muito usado em aplicações web.",
            "No projeto: roda dentro de um container; a API conecta via connection string.",
            "Banco criado automaticamente (MYSQL_DATABASE=MinhaPrimeiraApi).",
            "Tabelas: criadas pela aplicação via Entity Framework (EnsureCreated) na primeira execução.",
            "Dados persistem no volume Docker (mysql_data) mesmo após docker compose down."
        ],
        notes="Porta 3306 exposta para acesso externo (cliente MySQL, etc.)."
    )

    add_content_slide(
        prs,
        "Ferramenta 3: Entity Framework Core",
        [
            "EF Core: ORM (Object-Relational Mapper) da Microsoft para .NET.",
            "Permite trabalhar com banco usando classes C# (entidades) em vez de SQL direto.",
            "DbContext: representa a sessão com o banco; expõe DbSet<T> para cada tabela.",
            "Provedor MySQL: Pomelo.EntityFrameworkCore.MySql conecta o EF ao MySQL.",
            "No projeto: AppDbContext define Produtos e Categorias; EnsureCreated cria o banco e as tabelas."
        ],
        notes="Mencione que em produção é comum usar Migrations em vez de EnsureCreated."
    )

    add_content_slide(
        prs,
        "Fluxo da aplicação",
        [
            "Cliente HTTP (navegador, Postman, .http) envia requisição para a API.",
            "Controller (ex.: ProdutosController) recebe e usa o AppDbContext.",
            "DbContext executa as operações no MySQL (SELECT, INSERT, UPDATE, DELETE).",
            "Resposta JSON é devolvida ao cliente.",
            "Com Docker: API e MySQL na mesma rede; API usa hostname 'mysql' na connection string."
        ],
        notes="Desenhe na lousa: Cliente -> API (Controller -> DbContext) -> MySQL."
    )

    add_content_slide(
        prs,
        "Passo a passo: Pré-requisitos",
        [
            "Ter o Docker Desktop instalado e em execução (para rodar com Docker).",
            "Ou: ter .NET 9 SDK instalado (para rodar só a API com dotnet run).",
            "Ter o código do projeto (clone ou cópia da pasta MinhaPrimeiraApi).",
            "Abrir a pasta do projeto no terminal (PowerShell ou CMD)."
        ],
        notes="Docker Desktop: docker.com/products/docker-desktop"
    )

    add_content_slide(
        prs,
        "Passo a passo: Executar com Docker",
        [
            "1. Abra o terminal na pasta do projeto (onde está o docker-compose.yml).",
            "2. Execute: docker compose up --build -d",
            "   (--build reconstrói a imagem da API; -d roda em segundo plano).",
            "3. Aguarde o MySQL ficar saudável e a API subir (alguns segundos).",
            "4. A API estará em: http://localhost:8080",
            "   Ex.: http://localhost:8080/api/Produtos"
        ],
        notes="Na primeira vez, o download das imagens pode demorar."
    )

    add_content_slide(
        prs,
        "Passo a passo: Testar a API",
        [
            "Navegador: acesse http://localhost:8080/api/Produtos (lista em JSON).",
            "Arquivo .http: abra MinhaPrimeiraApi.http e use localhost:8080; execute cada bloco.",
            "Postman/Insomnia: GET/POST/PUT/DELETE em /api/Produtos e /api/Categorias.",
            "Parar os containers: docker compose down (dados do MySQL ficam no volume)."
        ],
        notes="Mostre na prática um GET e um POST no navegador ou Postman."
    )

    add_content_slide(
        prs,
        "Endpoints disponíveis (resumo)",
        [
            "Produtos: GET/POST /api/Produtos — GET/PUT/DELETE /api/Produtos/{id}",
            "Categorias: GET/POST /api/Categorias — GET/PUT/PATCH/DELETE /api/Categorias/{id}",
            "Ex.: POST /api/Produtos com body {\"nome\":\"...\", \"preco\": 100} cria um produto.",
            "OpenAPI (Development): http://localhost:8080/openapi/v1.json"
        ],
        notes="Ressalte a diferença entre PUT (substituição) e PATCH (parcial) em Categorias."
    )

    add_content_slide(
        prs,
        "Estrutura do projeto",
        [
            "Controllers/: ProdutosController, CategoriasController (rotas da API).",
            "Data/AppDbContext.cs: contexto EF Core (DbSet Produtos e Categorias).",
            "entities/: classes Produto e Categoria (tabelas no banco).",
            "DTOs/: objetos de entrada/saída (ex.: CreateCategoriaDto, CategoriaDto).",
            "Services/: CategoriaService (regras de negócio usando o DbContext).",
            "Dockerfile e docker-compose.yml: build e orquestração com MySQL."
        ],
        notes="Abra a árvore do projeto na IDE para mostrar aos alunos."
    )

    add_content_slide(
        prs,
        "Executar sem Docker (opcional)",
        [
            "Subir só o MySQL: docker compose up mysql -d",
            "Na pasta do projeto: dotnet run",
            "API em http://localhost:5155 (veja Properties/launchSettings.json).",
            "Connection string em appsettings.json: localhost, porta 3306, banco MinhaPrimeiraApi."
        ],
        notes="Útil quando o aluno quer debugar a API no Visual Studio sem container."
    )

    add_content_slide(
        prs,
        "Referências e dúvidas",
        [
            "Documentação Docker: docs.docker.com",
            "Entity Framework Core: learn.microsoft.com/ef/core",
            "MySQL: dev.mysql.com/doc",
            "Projeto: leia o README.md na raiz do repositório.",
            "Dúvidas: discutir em sala ou nos fóruns da disciplina."
        ],
        notes="Deixe tempo para perguntas e indique o README para consulta depois."
    )

    prs.save(OUTPUT_FILE)
    print(f"Apresentação gerada: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
